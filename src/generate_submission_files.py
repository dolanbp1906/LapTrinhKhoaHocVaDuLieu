"""Sinh báo cáo Word 20–30 trang và slide PPTX từ kết quả pipeline."""
from __future__ import annotations

import json
from pathlib import Path

import pandas as pd
from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PInches
from pptx.util import Pt as PPt

ROOT = Path(__file__).resolve().parents[1]
REPORTS = ROOT / "reports"
FIGURES = REPORTS / "figures"
PROCESSED = ROOT / "data" / "processed"
SUMMARIES = PROCESSED / "summaries"


def add_page_number(paragraph) -> None:
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = paragraph.add_run()
    fld_char1 = OxmlElement("w:fldChar")
    fld_char1.set(qn("w:fldCharType"), "begin")
    instr_text = OxmlElement("w:instrText")
    instr_text.set(qn("xml:space"), "preserve")
    instr_text.text = "PAGE"
    fld_char2 = OxmlElement("w:fldChar")
    fld_char2.set(qn("w:fldCharType"), "end")
    run._r.append(fld_char1)
    run._r.append(instr_text)
    run._r.append(fld_char2)


def set_cell_text(cell, text: object, bold: bool = False) -> None:
    cell.text = ""
    p = cell.paragraphs[0]
    run = p.add_run(str(text))
    run.bold = bold
    run.font.size = Pt(9)


def add_dataframe_table(doc: Document, df: pd.DataFrame, max_rows: int = 12) -> None:
    shown = df.head(max_rows).copy()
    table = doc.add_table(rows=1, cols=len(shown.columns))
    table.style = "Table Grid"
    for idx, col in enumerate(shown.columns):
        set_cell_text(table.rows[0].cells[idx], col, bold=True)
    for row in shown.itertuples(index=False):
        cells = table.add_row().cells
        for idx, value in enumerate(row):
            if isinstance(value, float):
                value = f"{value:,.2f}"
            set_cell_text(cells[idx], value)


def page_title(doc: Document, title: str, subtitle: str | None = None) -> None:
    doc.add_page_break()
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(title)
    run.bold = True
    run.font.size = Pt(20)
    run.font.color.rgb = RGBColor(31, 78, 121)
    if subtitle:
        p2 = doc.add_paragraph()
        p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r2 = p2.add_run(subtitle)
        r2.italic = True
        r2.font.size = Pt(11)


def add_body(doc: Document, paragraphs: list[str]) -> None:
    for text in paragraphs:
        p = doc.add_paragraph(text)
        p.paragraph_format.space_after = Pt(7)
        p.paragraph_format.line_spacing = 1.15


def add_figure(doc: Document, filename: str, caption: str, width: float = 5.8) -> None:
    path = FIGURES / filename
    if path.exists():
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run().add_picture(str(path), width=Inches(width))
        cap = doc.add_paragraph(caption)
        cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
        cap.runs[0].italic = True
        cap.runs[0].font.size = Pt(9)


def build_report() -> Path:
    fact = pd.read_csv(PROCESSED / "fact_sales.csv")
    products = pd.read_csv(PROCESSED / "products_final.csv")
    customers = pd.read_csv(PROCESSED / "customers.csv")
    orders = pd.read_csv(PROCESSED / "orders.csv")
    details = pd.read_csv(PROCESSED / "order_details.csv")
    inventory = pd.read_csv(PROCESSED / "inventory_transactions.csv")
    rfm = pd.read_csv(PROCESSED / "rfm_segments.csv")
    metrics = json.loads((REPORTS / "07_ml_metrics.json").read_text(encoding="utf-8"))
    numpy_evidence = json.loads(
        (REPORTS / "03_numpy_dictionary_evidence.json").read_text(encoding="utf-8")
    )

    doc = Document()
    sec = doc.sections[0]
    sec.top_margin = Inches(0.75)
    sec.bottom_margin = Inches(0.7)
    sec.left_margin = Inches(0.85)
    sec.right_margin = Inches(0.75)
    add_page_number(sec.footer.paragraphs[0])
    styles = doc.styles
    styles["Normal"].font.name = "Times New Roman"
    styles["Normal"].font.size = Pt(12)
    for style_name in ["Title", "Heading 1", "Heading 2", "Heading 3"]:
        styles[style_name].font.name = "Times New Roman"

    # Trang 1 — bìa
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("BÁO CÁO CHUYÊN ĐỀ CUỐI KỲ")
    r.bold = True
    r.font.size = Pt(18)
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("MÔN LẬP TRÌNH CHO KHOA HỌC DỮ LIỆU")
    r.bold = True
    r.font.size = Pt(16)
    doc.add_paragraph("\n")
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(
        "PHÂN TÍCH BÁN HÀNG, QUẢN LÝ TỒN KHO\nVÀ PHÂN NHÓM KHÁCH HÀNG"
    )
    r.bold = True
    r.font.size = Pt(22)
    r.font.color.rgb = RGBColor(31, 78, 121)
    doc.add_paragraph("\n")
    for line in [
        "Học viên: ........................................................",
        "Mã học viên: ...................................................",
        "Giảng viên: ......................................................",
        "Hình thức thực hiện: Cá nhân",
        "Phạm vi: Nhà sách và văn phòng phẩm",
    ]:
        p = doc.add_paragraph(line)
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph("\n\n")
    p = doc.add_paragraph("Tháng 08 năm 2026")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Trang 2
    page_title(doc, "LỜI CAM ĐOAN VÀ PHẠM VI SỬ DỤNG AI")
    add_body(
        doc,
        [
            "Báo cáo này trình bày dự án cá nhân về phân tích bán hàng, quản lý tồn kho và phân nhóm khách hàng. Dữ liệu giao dịch và khách hàng được mô phỏng bằng seed cố định; không chứa thông tin nhận diện khách hàng thật.",
            "AI được sử dụng để gợi ý cấu trúc, rà soát mã, giải thích lỗi, đề xuất kiểm thử và phản biện kết quả. Mọi đầu ra được kiểm tra lại bằng cách chạy pipeline, đối chiếu số dòng, khóa dữ liệu, chỉ số mô hình và biểu đồ. Nhật ký sử dụng AI được lưu tại logs/ai_usage_log.md.",
            "Nguồn trực tuyến chỉ được truy cập ở phạm vi công khai, có timeout và khoảng nghỉ. Dự án không vượt CAPTCHA, đăng nhập, token nội bộ hoặc cơ chế chống bot.",
            "Người thực hiện chịu trách nhiệm giải thích mã nguồn, dữ liệu, giả định, giới hạn và kết quả được trình bày trong báo cáo.",
        ],
    )
    doc.add_heading("Danh mục tệp minh chứng", level=2)
    for item in [
        "Mã nguồn theo lớp và module trong src/",
        "Bốn notebook bắt buộc trong notebooks/",
        "Dữ liệu gốc và dữ liệu sạch trong data/raw và data/processed",
        "Nhật ký lỗi, crawl và AI trong logs/",
        "Biểu đồ, metric và phân tích sai số trong reports/",
    ]:
        doc.add_paragraph(item, style="List Bullet")

    # Trang 3
    page_title(doc, "TÓM TẮT")
    add_body(
        doc,
        [
            f"Dự án xây dựng pipeline khoa học dữ liệu hoàn chỉnh cho cửa hàng nhà sách và văn phòng phẩm. Catalog cuối gồm {len(products)} sản phẩm; dữ liệu giao dịch gồm {len(orders):,} đơn, {len(details):,} dòng chi tiết, {len(customers):,} khách hàng và {len(inventory):,} giao dịch kho.",
            "Hệ thống nghiệp vụ được tổ chức theo hướng đối tượng với các lớp Product, Customer, Order, OrderItem, InventoryTransaction, InventoryManager và SalesManager. Giao dịch làm tồn kho âm bị từ chối; thao tác bán hàng có rollback khi một dòng không đủ hàng.",
            "Phần phân tích sử dụng NumPy, Pandas, Groupby, Pivot và trực quan hóa. Kết quả cho thấy doanh thu chịu ảnh hưởng mạnh bởi nhóm thiết bị văn phòng có đơn giá cao; kênh offline đóng góp lớn nhất và số sản phẩm dưới mức đặt lại tương đối cao.",
            "Bài toán học máy so sánh baseline với Linear/Logistic Regression và Decision Tree. Decision Tree đạt R² khoảng 0,79 cho giá trị đơn và F1 khoảng 0,98 cho cảnh báo tồn kho trên dữ liệu mô phỏng. RFM và K-Means tạo bốn phân khúc cùng chiến lược chăm sóc.",
            "Kết quả chỉ minh họa quy trình kỹ thuật, chưa phù hợp để tự động ra quyết định kinh doanh khi chưa kiểm chứng trên dữ liệu thực.",
        ],
    )
    doc.add_heading("Từ khóa", level=2)
    doc.add_paragraph(
        "OOP, quản lý tồn kho, bán hàng, NumPy, Pandas, EDA, hồi quy, phân lớp, RFM, K-Means."
    )

    # Trang 4
    page_title(doc, "MỤC LỤC VÀ CẤU TRÚC BÁO CÁO")
    add_body(
        doc,
        [
            "1. Giới thiệu và câu hỏi nghiên cứu",
            "2. Phạm vi, đạo đức và nguồn dữ liệu",
            "3. Thiết kế dữ liệu và kiến trúc pipeline",
            "4. OOP và nghiệp vụ kho",
            "5. Thu thập, làm sạch và kiểm tra chất lượng",
            "6. NumPy, dictionary lồng nhau, Merge/Groupby/Pivot",
            "7. EDA bán hàng, tồn kho và khách hàng",
            "8. Học máy dự đoán giá trị đơn",
            "9. Học máy cảnh báo tồn kho",
            "10. RFM, K-Means và chiến lược",
            "11. Giới hạn, tái lập và kết luận",
            "Phụ lục: cấu trúc tệp, cách chạy và câu hỏi bảo vệ.",
        ],
    )
    doc.add_paragraph(
        "Lưu ý: cập nhật trường họ tên, mã học viên và giảng viên trên trang bìa trước khi nộp.",
        style="Intense Quote",
    )

    # Trang 5
    page_title(doc, "1. GIỚI THIỆU VÀ CÂU HỎI NGHIÊN CỨU")
    add_body(
        doc,
        [
            "Cửa hàng cần theo dõi nhập, xuất, điều chỉnh tồn kho, doanh thu và hành vi khách hàng trong một quy trình thống nhất. Nếu chỉ lưu bảng giao dịch mà thiếu kiểm soát nghiệp vụ, tồn kho có thể âm hoặc số liệu báo cáo không nhất quán. Nếu chỉ xây dựng chương trình nghiệp vụ mà không phân tích dữ liệu, hệ thống khó hỗ trợ quyết định.",
            "Mục tiêu của chuyên đề là kết hợp lập trình Python hướng đối tượng với quy trình khoa học dữ liệu: thiết kế dữ liệu, thu thập, kiểm tra chất lượng, làm sạch, phân tích khám phá, học máy, phân cụm và báo cáo có thể tái lập.",
        ],
    )
    for question in [
        "Danh mục và kênh bán nào tạo doanh thu lớn nhất?",
        "Sản phẩm nào bán chạy, bán chậm và thường được mua cùng nhau?",
        "Thời điểm nào có doanh thu cao nhất?",
        "Khách hàng có thể chia thành những nhóm RFM nào?",
        "Có thể dự đoán giá trị đơn hàng với sai số bao nhiêu?",
        "Sản phẩm nào có nguy cơ xuống dưới mức cảnh báo?",
    ]:
        doc.add_paragraph(question, style="List Number")

    # Trang 6
    page_title(doc, "2. PHẠM VI, ĐẠO ĐỨC VÀ NGUỒN DỮ LIỆU")
    source_df = pd.DataFrame(
        [
            ["Giảng viên", 60, "Dữ liệu sản phẩm khởi đầu", "lecturer"],
            ["DummyJSON", 15, "API sản phẩm công khai", "public_api"],
            ["Books to Scrape", 15, "Website thực hành công khai", "public_website"],
            ["HTML mẫu", 60, "Đối chiếu parser, không tính SP mới", "practice"],
        ],
        columns=["Nguồn", "Số SP", "Vai trò", "source_type"],
    )
    add_dataframe_table(doc, source_df)
    add_body(
        doc,
        [
            "Dữ liệu sản phẩm công khai được bổ sung ở phạm vi nhỏ, có metadata URL và thời điểm thu thập. Giá USD được nhân 25.000 và giá GBP được nhân 33.000 để quy đổi sang VND. Tỷ giá cố định là giả định học tập, không phải tỷ giá giao dịch.",
            "Các trường product_id mới, popularity_weight, paired_product_id và một phần tồn kho được mô phỏng; cột simulated_fields ghi rõ theo từng bản ghi. Customers, orders, order_details và inventory_transactions là dữ liệu synthetic với seed=42.",
            "Không có tên, điện thoại, email hoặc địa chỉ chi tiết của khách hàng thật.",
        ],
    )

    # Trang 7
    page_title(doc, "3. THIẾT KẾ DỮ LIỆU")
    schema_df = pd.DataFrame(
        [
            ["products", len(products), "product_id", "Danh mục, giá, tồn đầu, reorder"],
            ["customers", len(customers), "customer_id", "Thành phố, loại KH, ngày tham gia"],
            ["orders", len(orders), "order_id", "Ngày, thanh toán, kênh, giảm giá"],
            ["order_details", len(details), "order_id + product_id", "Số lượng, giá bán"],
            ["inventory_transactions", len(inventory), "transaction_id", "Loại giao dịch, trạng thái"],
        ],
        columns=["Bảng", "Số dòng", "Khóa", "Nội dung"],
    )
    add_dataframe_table(doc, schema_df)
    add_body(
        doc,
        [
            "Quan hệ chính: customers 1–n orders; orders 1–n order_details; products 1–n order_details; products 1–n inventory_transactions. Khi tạo fact_sales, các phép merge dùng validate='many_to_one' để chặn cardinality ngoài dự kiến.",
            "Bảng fact_sales giữ cấp độ dòng chi tiết đơn, bổ sung thuộc tính khách hàng, sản phẩm, tháng, doanh thu dòng, phần giảm giá phân bổ và doanh thu ròng.",
        ],
    )

    # Trang 8
    page_title(doc, "4. KIẾN TRÚC PIPELINE VÀ KHẢ NĂNG TÁI LẬP")
    add_body(
        doc,
        [
            "Pipeline được chia thành module theo nhiệm vụ thay vì đặt toàn bộ mã trong notebook. Các script buổi 4–8 lần lượt làm sạch/sinh dữ liệu, merge/tổng hợp, EDA, học máy và RFM. run_all.py điều phối toàn bộ và dừng ngay nếu một bước lỗi.",
            "Chế độ mặc định dùng dữ liệu raw đã lưu để tránh phụ thuộc mạng. Tùy chọn --with-crawl kích hoạt lại bước thu thập. Việc tách raw và processed cho phép truy vết thay đổi và không sửa trực tiếp dữ liệu nguồn.",
            "Các tham số ngẫu nhiên quan trọng đều cố định ở 42; ngày tham chiếu RFM cố định là 2026-01-01. Metric, case sai, bảng tổng hợp và hình được ghi thành tệp độc lập.",
        ],
    )
    for command in [
        "python src/run_all.py",
        "python src/run_all.py --with-crawl",
        "python -m streamlit run src/app_dashboard.py",
    ]:
        p = doc.add_paragraph(command)
        p.style = "Intense Quote"

    # Trang 9
    page_title(doc, "5. OOP VÀ NGHIỆP VỤ KHO")
    class_df = pd.DataFrame(
        [
            ["Product", "Thông tin và số lượng hiện tại"],
            ["Customer", "Thông tin khách hàng ẩn danh"],
            ["OrderItem", "Một dòng sản phẩm trong đơn"],
            ["Order", "Header đơn và danh sách dòng"],
            ["InventoryTransaction", "Trạng thái trước/sau giao dịch"],
            ["InventoryManager", "Nhập/xuất/điều chỉnh, log và kiểm tra"],
            ["SalesManager", "Tạo đơn, trừ kho và rollback"],
        ],
        columns=["Lớp", "Trách nhiệm"],
    )
    add_dataframe_table(doc, class_df)
    add_body(
        doc,
        [
            "InventoryManager là nguồn sự thật cho số lượng hiện tại. Mọi thao tác được kiểm tra đầu vào, ghi quantity_before, quantity_after, người thực hiện, trạng thái và ghi chú.",
            "Khi xuất vượt tồn, hệ thống tạo giao dịch rejected và ghi error_log thay vì làm tồn kho âm. Khi tạo đơn nhiều dòng, nếu một dòng thất bại thì các dòng đã trừ trước đó được hoàn tác.",
        ],
    )

    # Trang 10
    page_title(doc, "6. THU THẬP VÀ KIỂM TRA CHẤT LƯỢNG")
    add_body(
        doc,
        [
            "Bước thu thập có timeout, kiểm tra HTTP status, khoảng nghỉ và xử lý lỗi theo nguồn. HTML mẫu hai trang được parse bằng selector cụ thể rồi đối chiếu 60/60 sản phẩm với dữ liệu giảng viên.",
            "Dữ liệu DummyJSON và Books to Scrape được ánh xạ về schema chung. Metadata source_type, source_reference, source_url và collected_at được giữ lại. Nhật ký crawl mô tả số lượng và lỗi.",
            "Các kiểm tra chất lượng bao gồm số ô thiếu, trùng product_id, giá không dương, initial_quantity không lớn hơn reorder_level, kiểu dữ liệu và tính hợp lệ của paired_product_id.",
        ],
    )
    quality_df = products.groupby("source_type", as_index=False).agg(
        products=("product_id", "count"),
        avg_price=("unit_price", "mean"),
        categories=("category", "nunique"),
    )
    add_dataframe_table(doc, quality_df)

    # Trang 11
    page_title(doc, "7. LÀM SẠCH, HỢP NHẤT VÀ KIỂM TRA KHÓA")
    add_body(
        doc,
        [
            "Chuỗi xử lý chuẩn hóa khoảng trắng, kiểu số, đơn vị, danh mục và thông tin nguồn; loại khóa trùng; sửa reorder_level không hợp lệ theo quy tắc công bố. Khi trùng product_id, bản ghi giảng viên được ưu tiên.",
            "Sau merge, catalog có 90 khóa product_id duy nhất, unit_price dương và initial_quantity lớn hơn reorder_level. Khóa ngoại từ order_details và inventory_transactions sang products được kiểm tra.",
            f"Fact table có {len(fact):,} dòng, bằng số dòng order_details. Các bước merge không làm tăng số dòng bất thường.",
        ],
    )
    merge_check = pd.DataFrame(
        [
            ["products.product_id duy nhất", products["product_id"].is_unique],
            ["orders.order_id duy nhất", orders["order_id"].is_unique],
            ["order_details FK order hợp lệ", details["order_id"].isin(orders["order_id"]).all()],
            ["order_details FK product hợp lệ", details["product_id"].isin(products["product_id"]).all()],
            ["inventory FK product hợp lệ", inventory["product_id"].isin(products["product_id"]).all()],
        ],
        columns=["Kiểm tra", "Kết quả"],
    )
    add_dataframe_table(doc, merge_check)

    # Trang 12
    page_title(doc, "8. DICTIONARY LỒNG NHAU VÀ NUMPY")
    np_result = numpy_evidence["numpy_analysis"]
    nested_result = numpy_evidence["nested_dictionary_analysis"]
    add_body(
        doc,
        [
            "Một mẫu giao dịch được biểu diễn theo cấu trúc order → customer/order/items. Vòng lặp trên dictionary lồng nhau tính doanh thu đơn, chi tiêu khách hàng và tần suất cặp sản phẩm mà không phụ thuộc Groupby.",
            f"Mẫu gồm {nested_result['sample_orders']} đơn với doanh thu {nested_result['sample_revenue']:,.0f} VND. Kết quả chi tiết được ghi JSON để kiểm tra.",
            f"Ma trận NumPy doanh thu tháng × danh mục có kích thước {tuple(np_result['matrix_shape'])}. sum(axis=1) tính tổng theo tháng; sum(axis=0) tính tổng theo danh mục. Z-score được tính theo cột và trung bình sau chuẩn hóa xấp xỉ 0.",
            f"Giá trị tồn kho đầu kỳ được vector hóa bằng phép nhân hai mảng, đạt {np_result['inventory_value_total']:,.0f} VND.",
        ],
    )
    matrix = pd.read_csv(SUMMARIES / "numpy_revenue_matrix.csv")
    add_dataframe_table(doc, matrix, max_rows=8)

    # Trang 13
    page_title(doc, "9. MERGE, GROUPBY VÀ PIVOT")
    add_body(
        doc,
        [
            "Các bảng tổng hợp gồm doanh thu theo tháng, danh mục, thành phố, kênh, phương thức thanh toán, top sản phẩm và tồn kho. Pivot tháng × kênh hỗ trợ so sánh cơ cấu theo thời gian; pivot danh mục × thanh toán hỗ trợ phân tích hành vi.",
            "Doanh thu dòng bằng quantity × selling_price. Discount cấp đơn được phân bổ theo tỷ trọng doanh thu dòng để tránh tính lặp khi tổng hợp.",
            "Tồn kho hiện tại được tính theo thứ tự timestamp và transaction_id, với dấu dương cho nhập, dấu âm cho bán/xuất/hao hụt điều chỉnh.",
        ],
    )
    add_dataframe_table(doc, pd.read_csv(SUMMARIES / "revenue_by_month.csv"))

    # Trang 14
    page_title(doc, "10. EDA — DOANH THU THEO THỜI GIAN")
    add_figure(doc, "01_revenue_by_month.png", "Hình 1. Doanh thu ròng theo tháng")
    add_body(
        doc,
        [
            "Doanh thu biến động theo tháng và đạt đỉnh vào 2025-04. Mức biến động phản ánh đồng thời số đơn, số lượng và cơ cấu mặt hàng có giá cao.",
            "Do dữ liệu giao dịch là synthetic, không nên diễn giải đỉnh tháng như quan hệ nhân quả hoặc mùa vụ thật. Kết quả phù hợp để kiểm tra kỹ thuật tổng hợp theo thời gian.",
        ],
    )

    # Trang 15
    page_title(doc, "11. EDA — DANH MỤC VÀ SẢN PHẨM")
    add_figure(doc, "02_revenue_by_category.png", "Hình 2. Doanh thu theo danh mục")
    add_figure(doc, "03_top_slow_products.png", "Hình 3. Sản phẩm bán chạy và bán chậm", 5.2)
    add_body(
        doc,
        [
            "Thiết bị văn phòng dẫn đầu doanh thu vì catalog bổ sung có laptop/tablet đơn giá cao. Kết quả bị chi phối bởi cơ cấu nguồn và tỷ giá cố định.",
            "Sản phẩm bán chạy được đo theo tổng số lượng; sản phẩm doanh thu cao được đo theo net_revenue. Hai tiêu chí có thể cho thứ hạng khác nhau.",
        ],
    )

    # Trang 16
    page_title(doc, "12. EDA — KÊNH BÁN VÀ THÀNH PHỐ")
    add_figure(doc, "04_revenue_by_channel.png", "Hình 4. Cơ cấu doanh thu theo kênh")
    add_figure(doc, "09_revenue_by_city.png", "Hình 5. Doanh thu theo thành phố", 5.0)
    add_body(
        doc,
        [
            "Kênh offline chiếm khoảng 42% doanh thu và TP. Hồ Chí Minh dẫn đầu trong dữ liệu mô phỏng. Các phân bố được sinh theo quy tắc có kiểm soát, không đại diện thị phần thật.",
        ],
    )

    # Trang 17
    page_title(doc, "13. EDA — GIÁ TRỊ ĐƠN VÀ CẶP SẢN PHẨM")
    add_figure(doc, "05_order_value_dist.png", "Hình 6. Phân bố giá trị đơn hàng")
    add_figure(doc, "06_pair_heatmap.png", "Hình 7. Ma trận đồng mua sản phẩm", 5.0)
    add_body(
        doc,
        [
            "Giá trị đơn lệch phải và có các ngoại lệ lớn do thiết bị giá cao. Điều này giải thích vì sao MAE, RMSE và R² cần được xem đồng thời.",
            "Cặp mua cùng được tạo từ các sản phẩm duy nhất trong cùng đơn; ma trận đối xứng và đường chéo bị đặt bằng 0.",
        ],
    )

    # Trang 18
    page_title(doc, "14. EDA — TỒN KHO VÀ RFM THÔ")
    add_figure(doc, "07_stock_vs_reorder.png", "Hình 8. Tồn kho hiện tại so với reorder level")
    add_figure(doc, "08_rfm_raw_dist.png", "Hình 9. Phân bố RFM thô", 5.0)
    add_body(
        doc,
        [
            "Có 77/90 sản phẩm bằng hoặc dưới reorder_level trong kịch bản mô phỏng. Tỷ lệ cao tạo bài toán phân lớp mất cân bằng và không phản ánh tỷ lệ thực tế của cửa hàng.",
            "Recency, Frequency và Monetary có thang đo khác nhau; Monetary lệch phải nên được log1p trước chuẩn hóa và phân cụm.",
        ],
    )

    # Trang 19
    page_title(doc, "15. HỌC MÁY — DỰ ĐOÁN GIÁ TRỊ ĐƠN")
    reg_metrics = pd.DataFrame(metrics["order_value"]["metrics"]).T.reset_index()
    reg_metrics = reg_metrics.rename(columns={"index": "Model"})
    add_dataframe_table(doc, reg_metrics)
    add_body(
        doc,
        [
            "Dữ liệu được chia train/test trước khi fit. ColumnTransformer xử lý biến số và biến phân loại trong Pipeline. Baseline DummyRegressor dùng trung bình làm mốc.",
            "DecisionTreeRegressor cho MAE khoảng 10,8 triệu VND và R² khoảng 0,79, tốt hơn rõ rệt so với baseline và Linear Regression. Tuy nhiên độ chính xác phụ thuộc mạnh vào quy tắc sinh dữ liệu và sản phẩm giá cao.",
            "Mười trường hợp sai lớn nhất được lưu riêng; phần lớn liên quan đơn có avg_selling_price cao hoặc cấu trúc số lượng hiếm.",
        ],
    )

    # Trang 20
    page_title(doc, "16. HỌC MÁY — CẢNH BÁO TỒN KHO")
    clf_metrics = pd.DataFrame(metrics["stock_alert"]["metrics"]).T.reset_index()
    clf_metrics = clf_metrics.rename(columns={"index": "Model"})
    add_dataframe_table(doc, clf_metrics)
    add_body(
        doc,
        [
            "Nhãn low_stock bằng 1 khi current_quantity ≤ reorder_level. Đặc trưng chỉ dùng thông tin đầu kỳ và bán/nhập nửa đầu năm để hạn chế rò rỉ từ kết quả cuối kỳ.",
            "DecisionTreeClassifier đạt F1 khoảng 0,98; Logistic Regression cũng vượt baseline. Vì lớp dương chiếm khoảng 85%, Accuracy không đủ để đánh giá và phải xem Precision, Recall, F1.",
            "Tập chỉ có 90 sản phẩm nên kết quả có phương sai cao. Mô hình không nên tự động đặt hàng nếu chưa có dữ liệu lịch sử thực và kiểm thử theo thời gian.",
        ],
    )

    # Trang 21
    page_title(doc, "17. RFM VÀ K-MEANS")
    add_figure(doc, "11_rfm_clusters.png", "Hình 10. Các cụm khách hàng RFM")
    add_figure(doc, "12_rfm_segment_counts.png", "Hình 11. Quy mô phân khúc", 5.0)
    add_body(
        doc,
        [
            "Ngày tham chiếu cố định 2026-01-01. Recency tính số ngày từ lần mua gần nhất; Frequency là số đơn duy nhất; Monetary là tổng doanh thu ròng.",
            "Dữ liệu dùng recency_days, frequency và log1p(monetary), sau đó StandardScaler. Số cụm được lựa chọn bằng silhouette trong khoảng k=3–6; k=4 cho kết quả tốt nhất trong lần chạy hiện tại.",
        ],
    )

    # Trang 22
    page_title(doc, "18. CHIẾN LƯỢC THEO PHÂN KHÚC")
    segment_summary = rfm.groupby("segment", as_index=False).agg(
        customers=("customer_id", "count"),
        recency_mean=("recency_days", "mean"),
        frequency_mean=("frequency", "mean"),
        monetary_mean=("monetary", "mean"),
    )
    add_dataframe_table(doc, segment_summary)
    strategies = [
        ("Champions", "Ưu đãi VIP, early access, chăm sóc cá nhân hóa và giữ chân."),
        ("Loyal", "Tích điểm, upsell/cross-sell và gợi ý sản phẩm mua cùng."),
        ("Potential", "Coupon kích hoạt đơn tiếp theo và nội dung theo danh mục quan tâm."),
        ("At Risk", "Chiến dịch win-back có thời hạn, khảo sát lý do và giới hạn chi phí."),
    ]
    for name, strategy in strategies:
        p = doc.add_paragraph()
        p.add_run(f"{name}: ").bold = True
        p.add_run(strategy)
    add_body(
        doc,
        [
            "Tên phân khúc được gán sau khi xem trung bình cụm, không phải nhãn thật. Khi triển khai cần theo dõi tỷ lệ phản hồi, chi phí và quyền riêng tư.",
        ],
    )

    # Trang 23
    page_title(doc, "19. GIỚI HẠN, RỦI RO VÀ HƯỚNG PHÁT TRIỂN")
    limits = [
        "Giao dịch và khách hàng là synthetic nên không phản ánh hành vi thực.",
        "FX cố định và thiết bị giá cao làm doanh thu lệch mạnh.",
        "Tập tồn kho chỉ có 90 sản phẩm và mất cân bằng lớp.",
        "Chia train/test ngẫu nhiên phù hợp minh họa nhưng chưa phải đánh giá theo thời gian.",
        "RFM phụ thuộc ngày tham chiếu và quy tắc gán tên cụm.",
        "Nguồn web có thể thay đổi cấu trúc hoặc ngừng truy cập.",
    ]
    for item in limits:
        doc.add_paragraph(item, style="List Bullet")
    doc.add_heading("Hướng phát triển", level=2)
    for item in [
        "Thu thập lịch sử bán hàng thật đã ẩn danh và thiết kế đánh giá theo thời gian.",
        "Dự báo nhu cầu theo sản phẩm, thời gian giao hàng và chi phí thiếu hàng.",
        "Hiệu chỉnh ngưỡng reorder theo service level thay vì nhãn cố định.",
        "Theo dõi drift mô hình và hiệu quả chiến dịch theo phân khúc.",
        "Bổ sung kiểm thử tự động, data validation và version hóa dữ liệu.",
    ]:
        doc.add_paragraph(item, style="List Bullet")

    # Trang 24
    page_title(doc, "20. KẾT LUẬN")
    add_body(
        doc,
        [
            "Dự án đáp ứng chuỗi yêu cầu của Chuyên đề 6: OOP quản lý kho, File I/O, dictionary lồng nhau, NumPy, Pandas, Merge/Groupby/Pivot, EDA, mô hình giám sát, phân tích sai số và RFM/K-Means.",
            f"Quy mô dữ liệu vượt mức tối thiểu: {len(products)} sản phẩm, {len(customers)} khách hàng, {len(orders):,} đơn và {len(details):,} dòng chi tiết. Pipeline tạo lại các bảng, biểu đồ, metric và phân khúc bằng một lệnh.",
            "Kết quả học máy cho thấy Decision Tree phù hợp hơn baseline trong dữ liệu mô phỏng, nhưng không chứng minh khả năng vận hành trên dữ liệu thật. Giá trị chính của dự án là quy trình có kiểm soát, minh bạch về nguồn và giả định.",
        ],
    )
    doc.add_heading("Lệnh chạy", level=2)
    doc.add_paragraph("python src/run_all.py", style="Intense Quote")
    doc.add_paragraph("python -m streamlit run src/app_dashboard.py", style="Intense Quote")

    # Trang 25
    page_title(doc, "PHỤ LỤC A — CẤU TRÚC SẢN PHẨM NỘP")
    for item in [
        "notebooks/01_problem_and_data.ipynb",
        "notebooks/02_collection_and_cleaning.ipynb",
        "notebooks/03_eda.ipynb",
        "notebooks/04_machine_learning.ipynb",
        "src/models.py, inventory_manager.py, sales_manager.py",
        "src/crawl_*.py, clean_products.py, merge_products.py, validate_products.py",
        "src/run_buoi4.py … run_buoi8_rfm.py, run_numpy_evidence.py, run_all.py",
        "data/raw/, data/processed/, logs/, reports/figures/",
        "reports/BAO_CAO_CHUYEN_DE_6.docx",
        "reports/SLIDE_CHUYEN_DE_6.pptx",
    ]:
        doc.add_paragraph(item, style="List Bullet")
    doc.add_heading("Kiểm tra trước khi nộp", level=2)
    for item in [
        "Điền họ tên, mã học viên và tên giảng viên.",
        "Mở Word, cập nhật mục lục nếu cần và kiểm tra ngắt trang.",
        "Kiểm tra font, hình, số liệu và các liên kết tệp.",
        "Chạy pipeline một lần cuối và lưu log.",
    ]:
        doc.add_paragraph(item, style="List Number")

    # Trang 26
    page_title(doc, "PHỤ LỤC B — CÂU HỎI BẢO VỆ GỢI Ý")
    questions = [
        ("Vì sao cần validate cardinality khi merge?", "Để phát hiện khóa trùng gây phình dòng và làm sai doanh thu."),
        ("Vì sao RMSE lớn hơn MAE?", "RMSE phạt mạnh các sai số lớn; dữ liệu có outlier thiết bị giá cao."),
        ("Vì sao không chỉ dùng Accuracy?", "Lớp low_stock mất cân bằng, cần xem Precision, Recall và F1."),
        ("Làm sao tránh rò rỉ tồn kho?", "Chỉ dùng đặc trưng có trước thời điểm dự báo, không dùng lượng bán cả kỳ."),
        ("Vì sao log Monetary?", "Monetary lệch phải; log1p giảm ảnh hưởng khách chi tiêu cực lớn."),
        ("Vì sao cố định reference_date?", "Để RFM tái lập và không đổi theo ngày chạy."),
        ("Điểm yếu lớn nhất?", "Dữ liệu synthetic và quy mô tồn kho nhỏ, chưa chứng minh hiệu quả thực tế."),
    ]
    for q, a in questions:
        p = doc.add_paragraph()
        p.add_run(f"Hỏi: {q}").bold = True
        doc.add_paragraph(f"Đáp: {a}")

    out = REPORTS / "BAO_CAO_CHUYEN_DE_6.docx"
    doc.save(out)
    return out


NAVY = PRGBColor(31, 78, 121)
INK = PRGBColor(15, 32, 51)
INK_SOFT = PRGBColor(38, 62, 88)
ACCENT = PRGBColor(0, 150, 158)
ACCENT_DARK = PRGBColor(0, 110, 118)
SURFACE = PRGBColor(244, 247, 250)
BORDER = PRGBColor(216, 224, 233)
WHITE = PRGBColor(255, 255, 255)
MUTED = PRGBColor(104, 122, 142)
FONT = "Segoe UI"

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.72


def style_text(
    frame,
    lines: list[tuple[str, float, bool, PRGBColor]],
    align=PP_ALIGN.LEFT,
    line_spacing: float | None = None,
    space_after: float = 6,
) -> None:
    frame.word_wrap = True
    for idx, (text, size, bold, color) in enumerate(lines):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = text
        p.alignment = align
        p.space_after = PPt(space_after)
        if line_spacing:
            p.line_spacing = line_spacing
        font = p.font
        font.size = PPt(size)
        font.bold = bold
        font.color.rgb = color
        font.name = FONT


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: PRGBColor | None,
    line: PRGBColor | None = None,
    shape=MSO_SHAPE.RECTANGLE,
    radius: float | None = None,
):
    shp = slide.shapes.add_shape(shape, PInches(x), PInches(y), PInches(w), PInches(h))
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    if radius is not None and shp.adjustments:
        shp.adjustments[0] = radius
    shp.text_frame.text = ""
    return shp


def add_text(slide, x: float, y: float, w: float, h: float):
    box = slide.shapes.add_textbox(PInches(x), PInches(y), PInches(w), PInches(h))
    box.text_frame.word_wrap = True
    return box.text_frame


def content_slide(prs, kicker: str, title: str, number: int):
    """Slide nội dung có thanh nhấn, tiêu đề và chân trang thống nhất."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, SURFACE)
    add_rect(slide, 0, 0, SLIDE_W, 0.11, ACCENT)

    frame = add_text(slide, MARGIN, 0.34, 11.0, 0.9)
    style_text(
        frame,
        [
            (kicker.upper(), 11.5, True, ACCENT_DARK),
            (title, 27, True, INK),
        ],
        space_after=2,
    )

    add_rect(slide, MARGIN, SLIDE_H - 0.62, SLIDE_W - 2 * MARGIN, 0.012, BORDER)
    foot = add_text(slide, MARGIN, SLIDE_H - 0.56, 9.0, 0.34)
    style_text(foot, [("Chuyên đề 6 — Phân tích bán hàng, tồn kho và phân nhóm khách hàng", 9.5, False, MUTED)])
    num = add_text(slide, SLIDE_W - MARGIN - 1.0, SLIDE_H - 0.56, 1.0, 0.34)
    style_text(num, [(f"{number:02d}", 9.5, True, MUTED)], align=PP_ALIGN.RIGHT)
    return slide


def add_bullets(slide, bullets: list[str], x: float, y: float, w: float, size: float = 16.5) -> None:
    """Danh sách gạch đầu dòng dùng ô vuông nhấn màu thay cho bullet mặc định."""
    step = 0.62 if size >= 16 else 0.54
    for idx, text in enumerate(bullets):
        top = y + idx * step
        add_rect(slide, x, top + 0.09, 0.11, 0.11, ACCENT)
        frame = add_text(slide, x + 0.28, top - 0.02, w - 0.28, step)
        style_text(frame, [(text, size, False, INK_SOFT)], line_spacing=1.05, space_after=0)


def add_kpi_cards(slide, cards: list[tuple[str, str]], y: float, height: float = 1.15) -> None:
    total_w = SLIDE_W - 2 * MARGIN
    gap = 0.22
    card_w = (total_w - gap * (len(cards) - 1)) / len(cards)
    for idx, (value, label) in enumerate(cards):
        x = MARGIN + idx * (card_w + gap)
        add_rect(slide, x, y, card_w, height, WHITE, BORDER, MSO_SHAPE.ROUNDED_RECTANGLE, 0.08)
        add_rect(slide, x, y, 0.07, height, ACCENT)
        frame = add_text(slide, x + 0.26, y + 0.14, card_w - 0.42, height - 0.24)
        style_text(
            frame,
            [(value, 23, True, INK), (label, 10.5, False, MUTED)],
            space_after=1,
        )


def add_panel(slide, x: float, y: float, w: float, h: float, title: str | None = None):
    add_rect(slide, x, y, w, h, WHITE, BORDER, MSO_SHAPE.ROUNDED_RECTANGLE, 0.03)
    if title:
        frame = add_text(slide, x + 0.28, y + 0.16, w - 0.5, 0.36)
        style_text(frame, [(title.upper(), 11, True, ACCENT_DARK)])


def add_figure_card(
    slide,
    filename: str,
    x: float,
    y: float,
    w: float,
    h: float,
    caption: str | None = None,
) -> None:
    """Đặt hình trong khung trắng, giữ đúng tỷ lệ ảnh gốc."""
    path = FIGURES / filename
    if not path.exists():
        return
    add_panel(slide, x, y, w, h)
    pad = 0.16
    caption_h = 0.34 if caption else 0.0
    box_w = w - 2 * pad
    box_h = h - 2 * pad - caption_h
    with Image.open(path) as img:
        ratio = img.width / img.height
    draw_w = box_w
    draw_h = draw_w / ratio
    if draw_h > box_h:
        draw_h = box_h
        draw_w = draw_h * ratio
    left = x + pad + (box_w - draw_w) / 2
    top = y + pad + (box_h - draw_h) / 2
    slide.shapes.add_picture(str(path), PInches(left), PInches(top), width=PInches(draw_w))
    if caption:
        frame = add_text(slide, x + pad, y + h - pad - caption_h + 0.02, box_w, caption_h)
        style_text(frame, [(caption, 10.5, False, MUTED)], align=PP_ALIGN.CENTER)


def add_metric_table(
    slide,
    headers: list[str],
    rows: list[list[str]],
    x: float,
    y: float,
    w: float,
    highlight_last: bool = True,
) -> None:
    row_h = 0.42
    shape = slide.shapes.add_table(
        len(rows) + 1, len(headers), PInches(x), PInches(y), PInches(w), PInches(row_h * (len(rows) + 1))
    )
    table = shape.table
    table.first_row = True
    if len(headers) > 1:
        first_w = w * 0.44
        rest_w = (w - first_w) / (len(headers) - 1)
        table.columns[0].width = PInches(first_w)
        for idx in range(1, len(headers)):
            table.columns[idx].width = PInches(rest_w)
    for idx, header in enumerate(headers):
        cell = table.cell(0, idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = INK
        cell.margin_left = PInches(0.12)
        style_text(
            cell.text_frame,
            [(header, 11.5, True, WHITE)],
            align=PP_ALIGN.LEFT if idx == 0 else PP_ALIGN.RIGHT,
            space_after=0,
        )
    for r, row in enumerate(rows, start=1):
        is_best = highlight_last and r == len(rows)
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PRGBColor(233, 246, 246) if is_best else WHITE
            cell.margin_left = PInches(0.12)
            style_text(
                cell.text_frame,
                [(value, 11.5, is_best, INK if is_best else INK_SOFT)],
                align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT,
                space_after=0,
            )


def section_slide(prs, label: str, title: str, note: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, INK)
    add_rect(slide, 0, 0, 0.16, SLIDE_H, ACCENT)
    frame = add_text(slide, 1.5, 2.7, 10.5, 2.2)
    style_text(
        frame,
        [
            (label.upper(), 13, True, ACCENT),
            (title, 40, True, WHITE),
            (note, 14, False, PRGBColor(178, 195, 210)),
        ],
        space_after=8,
    )


def build_slides() -> Path:
    metrics = json.loads((REPORTS / "07_ml_metrics.json").read_text(encoding="utf-8"))
    rfm = pd.read_csv(PROCESSED / "rfm_segments.csv")
    reg = metrics["order_value"]["metrics"]
    clf = metrics["stock_alert"]["metrics"]

    prs = Presentation()
    prs.slide_width = PInches(SLIDE_W)
    prs.slide_height = PInches(SLIDE_H)
    blank = prs.slide_layouts[6]

    # --- 1. Trang bìa ---
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, INK)
    add_rect(slide, 0, 0, SLIDE_W, 0.16, ACCENT)
    add_rect(slide, 7.9, 0, 5.44, SLIDE_H, INK_SOFT)
    add_rect(slide, 7.9, 0, 0.05, SLIDE_H, ACCENT)

    frame = add_text(slide, 1.0, 1.55, 6.6, 0.5)
    style_text(frame, [("CHUYÊN ĐỀ CUỐI KỲ · SỐ 06", 13, True, ACCENT)])
    frame = add_text(slide, 1.0, 2.05, 6.7, 2.5)
    style_text(
        frame,
        [
            ("Phân tích bán hàng,", 33, True, WHITE),
            ("quản lý tồn kho và", 33, True, WHITE),
            ("phân nhóm khách hàng", 33, True, WHITE),
        ],
        line_spacing=1.05,
        space_after=0,
    )
    add_rect(slide, 1.0, 4.72, 1.5, 0.045, ACCENT)
    frame = add_text(slide, 1.0, 4.95, 6.6, 1.2)
    style_text(
        frame,
        [
            ("Lập trình cho Khoa học dữ liệu · Nhà sách và văn phòng phẩm", 13, False, PRGBColor(178, 195, 210)),
            ("Học viên: ……………………………    MSSV: ……………………", 13, False, PRGBColor(178, 195, 210)),
            ("Giảng viên: ……………………………             Tháng 08/2026", 13, False, PRGBColor(178, 195, 210)),
        ],
        space_after=5,
    )

    highlights = [
        ("90", "sản phẩm sau hợp nhất"),
        ("1.200", "đơn hàng · 3.274 dòng"),
        ("R² 0,79", "dự đoán giá trị đơn"),
        ("F1 0,98", "cảnh báo tồn kho"),
    ]
    for idx, (value, label) in enumerate(highlights):
        top = 1.6 + idx * 1.12
        add_rect(slide, 8.5, top, 4.3, 0.92, INK, None, MSO_SHAPE.ROUNDED_RECTANGLE, 0.1)
        add_rect(slide, 8.5, top, 0.06, 0.92, ACCENT)
        frame = add_text(slide, 8.75, top + 0.11, 4.0, 0.7)
        style_text(frame, [(value, 20, True, WHITE), (label, 10.5, False, PRGBColor(150, 172, 192))], space_after=0)

    # --- 2. Nội dung trình bày ---
    slide = content_slide(prs, "Nội dung", "Trình tự trình bày", 2)
    agenda = [
        ("01", "Bài toán và câu hỏi nghiên cứu"),
        ("02", "Dữ liệu, nguồn và phạm vi"),
        ("03", "Kiến trúc dữ liệu & OOP kho"),
        ("04", "EDA bán hàng và tồn kho"),
        ("05", "Học máy: giá trị đơn & cảnh báo"),
        ("06", "RFM, K-Means và chiến lược"),
    ]
    for idx, (num, text) in enumerate(agenda):
        col = idx % 2
        row = idx // 2
        x = MARGIN + col * 6.05
        y = 1.72 + row * 1.35
        add_rect(slide, x, y, 5.75, 1.1, WHITE, BORDER, MSO_SHAPE.ROUNDED_RECTANGLE, 0.06)
        add_rect(slide, x + 0.22, y + 0.24, 0.62, 0.62, ACCENT, None, MSO_SHAPE.ROUNDED_RECTANGLE, 0.15)
        frame = add_text(slide, x + 0.22, y + 0.36, 0.62, 0.4)
        style_text(frame, [(num, 14, True, WHITE)], align=PP_ALIGN.CENTER, space_after=0)
        frame = add_text(slide, x + 1.02, y + 0.36, 4.5, 0.5)
        style_text(frame, [(text, 15, False, INK)], space_after=0)

    # --- 3. Bài toán ---
    slide = content_slide(prs, "Bối cảnh", "Bài toán và mục tiêu", 3)
    add_panel(slide, MARGIN, 1.55, 6.1, 4.85, "Mục tiêu hệ thống")
    add_bullets(
        slide,
        [
            "Nhập / xuất / điều chỉnh kho có nhật ký",
            "Chặn tồn kho âm, rollback khi thiếu hàng",
            "Phân tích doanh thu và cặp sản phẩm",
            "Dự đoán giá trị đơn hàng",
            "Cảnh báo sản phẩm dưới reorder",
            "Phân nhóm khách hàng RFM",
        ],
        MARGIN + 0.3,
        2.2,
        5.5,
        15.5,
    )
    add_panel(slide, 7.1, 1.55, 5.5, 4.85, "6 câu hỏi nghiên cứu")
    questions = [
        "Danh mục / kênh nào doanh thu lớn nhất?",
        "Sản phẩm bán chạy, chậm, mua cùng nhau?",
        "Thời điểm nào doanh thu cao nhất?",
        "Khách hàng chia thành nhóm RFM nào?",
        "Dự đoán giá trị đơn sai số bao nhiêu?",
        "Sản phẩm nào nguy cơ dưới cảnh báo?",
    ]
    for idx, text in enumerate(questions):
        top = 2.18 + idx * 0.66
        add_rect(slide, 7.4, top, 0.34, 0.34, SURFACE, None, MSO_SHAPE.OVAL)
        frame = add_text(slide, 7.4, top + 0.03, 0.34, 0.3)
        style_text(frame, [(str(idx + 1), 11, True, ACCENT_DARK)], align=PP_ALIGN.CENTER, space_after=0)
        frame = add_text(slide, 7.85, top + 0.02, 4.5, 0.5)
        style_text(frame, [(text, 13.5, False, INK_SOFT)], line_spacing=1.0, space_after=0)

    # --- 4. Dữ liệu ---
    slide = content_slide(prs, "Dữ liệu", "Nguồn dữ liệu và quy mô", 4)
    add_kpi_cards(
        slide,
        [("90", "SẢN PHẨM"), ("220", "KHÁCH HÀNG"), ("1.200", "ĐƠN HÀNG"), ("3.274", "DÒNG CHI TIẾT"), ("3.396", "GIAO DỊCH KHO")],
        1.5,
    )
    add_metric_table(
        slide,
        ["Nguồn", "Số SP", "Loại"],
        [
            ["Giảng viên cung cấp", "60", "lecturer"],
            ["DummyJSON Products API", "15", "public_api"],
            ["Books to Scrape", "15", "public_website"],
            ["HTML mẫu (không tính SP mới)", "60", "practice"],
        ],
        MARGIN,
        3.05,
        6.1,
        highlight_last=False,
    )
    add_panel(slide, 7.1, 3.05, 5.5, 3.35, "Minh bạch dữ liệu")
    add_bullets(
        slide,
        [
            "Giá USD × 25.000, GBP × 33.000",
            "Cột simulated_fields ghi trường mô phỏng",
            "Giao dịch synthetic, seed = 42",
            "Không có dữ liệu khách hàng thật",
            "Không vượt CAPTCHA hoặc chống bot",
        ],
        7.4,
        3.7,
        5.0,
        13.5,
    )

    # --- 5. Kiến trúc + OOP ---
    slide = content_slide(prs, "Kiến trúc", "Pipeline dữ liệu và OOP quản lý kho", 5)
    steps = ["RAW", "CLEAN", "FACT", "EDA", "ML", "RFM"]
    for idx, step in enumerate(steps):
        x = MARGIN + idx * 2.03
        add_rect(slide, x, 1.55, 1.78, 0.62, WHITE, BORDER, MSO_SHAPE.ROUNDED_RECTANGLE, 0.12)
        frame = add_text(slide, x, 1.68, 1.78, 0.4)
        style_text(frame, [(step, 13, True, ACCENT_DARK)], align=PP_ALIGN.CENTER, space_after=0)
        if idx < len(steps) - 1:
            add_rect(slide, x + 1.83, 1.83, 0.14, 0.05, ACCENT)

    add_panel(slide, MARGIN, 2.45, 6.1, 3.95, "7 lớp theo đề cương")
    add_bullets(
        slide,
        [
            "Product · Customer · Order · OrderItem",
            "InventoryTransaction",
            "InventoryManager — nhập/xuất/điều chỉnh",
            "SalesManager — tạo đơn, trừ kho, rollback",
        ],
        MARGIN + 0.3,
        3.1,
        5.5,
        15,
    )
    add_panel(slide, 7.1, 2.45, 5.5, 3.95, "Kiểm soát nghiệp vụ")
    add_bullets(
        slide,
        [
            "Từ chối giao dịch làm âm kho",
            "Ghi inventory_log.csv + error_log.txt",
            "Lưu trạng thái trước / sau",
            "Merge validate many-to-one",
            "3.274 dòng fact, không phình dòng",
        ],
        7.4,
        3.1,
        5.0,
        14,
    )

    # --- 6. EDA doanh thu ---
    slide = content_slide(prs, "EDA", "Doanh thu theo thời gian và kênh bán", 6)
    add_figure_card(slide, "01_revenue_by_month.png", MARGIN, 1.5, 6.0, 3.85)
    add_figure_card(slide, "04_revenue_by_channel.png", 7.0, 1.5, 5.6, 3.85)
    add_kpi_cards(
        slide,
        [("2025-04", "THÁNG ĐỈNH DOANH THU"), ("~42%", "TỶ TRỌNG KÊNH OFFLINE"), ("TP. HCM", "THÀNH PHỐ DẪN ĐẦU")],
        5.55,
        0.82,
    )

    # --- 7. EDA sản phẩm & tồn kho ---
    slide = content_slide(prs, "EDA", "Sản phẩm, cặp mua cùng và tồn kho", 7)
    add_figure_card(slide, "07_stock_vs_reorder.png", MARGIN, 1.5, 6.0, 3.85)
    add_figure_card(slide, "06_pair_heatmap.png", 7.0, 1.5, 5.6, 3.85)
    add_kpi_cards(
        slide,
        [("77/90", "SP DƯỚI HOẶC BẰNG REORDER"), ("P0064–P0067", "CẶP ĐỒNG MUA MẠNH NHẤT"), ("Lệch phải", "PHÂN BỐ GIÁ TRỊ ĐƠN")],
        5.55,
        0.82,
    )

    # --- 8. ML giá trị đơn ---
    slide = content_slide(prs, "Học máy", "Dự đoán giá trị đơn hàng (hồi quy)", 8)
    add_metric_table(
        slide,
        ["Mô hình", "MAE (triệu)", "RMSE (triệu)", "R²"],
        [
            ["DummyRegressor (baseline)", f"{reg['DummyRegressor']['MAE']/1e6:,.1f}", f"{reg['DummyRegressor']['RMSE']/1e6:,.1f}", f"{reg['DummyRegressor']['R2']:.3f}"],
            ["LinearRegression", f"{reg['LinearRegression']['MAE']/1e6:,.1f}", f"{reg['LinearRegression']['RMSE']/1e6:,.1f}", f"{reg['LinearRegression']['R2']:.3f}"],
            ["DecisionTreeRegressor", f"{reg['DecisionTreeRegressor']['MAE']/1e6:,.1f}", f"{reg['DecisionTreeRegressor']['RMSE']/1e6:,.1f}", f"{reg['DecisionTreeRegressor']['R2']:.3f}"],
        ],
        MARGIN,
        1.6,
        7.4,
    )
    add_panel(slide, 8.4, 1.6, 4.2, 2.2, "Mô hình tốt nhất")
    frame = add_text(slide, 8.7, 2.15, 3.7, 1.3)
    style_text(
        frame,
        [("Decision Tree", 19, True, INK), ("MAE ≈ 10,8 triệu · R² ≈ 0,79", 12.5, False, MUTED)],
        space_after=3,
    )
    add_panel(slide, MARGIN, 4.05, 11.9, 2.35, "Thiết lập và phân tích sai số")
    add_bullets(
        slide,
        [
            "Chia train/test 900/300 trước khi fit; dùng Pipeline + ColumnTransformer",
            "Đặc trưng: số dòng, số lượng, giá TB, discount, tháng, thứ, kênh, thành phố, loại KH",
            "10 case sai lớn nhất: đơn có avg_selling_price rất cao hoặc cấu trúc số lượng hiếm",
        ],
        MARGIN + 0.3,
        4.7,
        11.2,
        14.5,
    )

    # --- 9. ML cảnh báo tồn kho ---
    slide = content_slide(prs, "Học máy", "Phân lớp cảnh báo tồn kho thấp", 9)
    add_metric_table(
        slide,
        ["Mô hình", "Accuracy", "Precision", "Recall", "F1"],
        [
            ["DummyClassifier (baseline)", f"{clf['DummyClassifier']['Accuracy']:.3f}", f"{clf['DummyClassifier']['Precision']:.3f}", f"{clf['DummyClassifier']['Recall']:.3f}", f"{clf['DummyClassifier']['F1']:.3f}"],
            ["LogisticRegression", f"{clf['LogisticRegression']['Accuracy']:.3f}", f"{clf['LogisticRegression']['Precision']:.3f}", f"{clf['LogisticRegression']['Recall']:.3f}", f"{clf['LogisticRegression']['F1']:.3f}"],
            ["DecisionTreeClassifier", f"{clf['DecisionTreeClassifier']['Accuracy']:.3f}", f"{clf['DecisionTreeClassifier']['Precision']:.3f}", f"{clf['DecisionTreeClassifier']['Recall']:.3f}", f"{clf['DecisionTreeClassifier']['F1']:.3f}"],
        ],
        MARGIN,
        1.6,
        8.3,
    )
    add_panel(slide, 9.3, 1.6, 3.3, 2.2, "Kết quả")
    frame = add_text(slide, 9.6, 2.15, 2.8, 1.3)
    style_text(frame, [("F1 ≈ 0,98", 19, True, INK), ("1 dự đoán sai trên test", 12.5, False, MUTED)], space_after=3)
    add_panel(slide, MARGIN, 4.05, 11.9, 2.35, "Chống rò rỉ và lưu ý đánh giá")
    add_bullets(
        slide,
        [
            "Nhãn: current_quantity ≤ reorder_level",
            "Chỉ dùng đặc trưng đầu kỳ + bán/nhập nửa đầu năm 2025 để tránh rò rỉ",
            "Lớp dương ~85% → ưu tiên Precision / Recall / F1 thay vì Accuracy",
        ],
        MARGIN + 0.3,
        4.7,
        11.2,
        14.5,
    )

    # --- 10. RFM ---
    slide = content_slide(prs, "Phân cụm", "RFM và K-Means phân nhóm khách hàng", 10)
    add_figure_card(slide, "11_rfm_clusters.png", MARGIN, 1.5, 5.9, 3.55)
    order = ["Champions", "Loyal", "Potential", "At Risk"]
    summary = rfm.groupby("segment").agg(n=("customer_id", "count"), r=("recency_days", "mean"), f=("frequency", "mean"))
    rows = []
    for name in order:
        if name in summary.index:
            row = summary.loc[name]
            rows.append([name, f"{int(row['n'])}", f"{row['r']:.0f}", f"{row['f']:.1f}"])
    add_metric_table(slide, ["Phân khúc", "Số KH", "Recency", "Frequency"], rows, 6.9, 1.6, 5.7, highlight_last=False)
    add_panel(slide, 6.9, 3.85, 5.7, 2.55, "Chiến lược")
    add_bullets(
        slide,
        [
            "Champions — ưu đãi VIP, giữ chân",
            "Loyal — upsell, tích điểm",
            "Potential — coupon kích hoạt",
            "At Risk — win-back có thời hạn",
        ],
        7.2,
        4.42,
        5.1,
        13,
    )
    add_panel(slide, MARGIN, 5.2, 5.9, 1.2, "Thiết lập")
    frame = add_text(slide, MARGIN + 0.3, 5.74, 5.4, 0.5)
    style_text(
        frame,
        [("reference_date = 2026-01-01 · k = 4 · log1p(Monetary)", 12.5, False, INK_SOFT)],
        space_after=0,
    )

    # --- 11. Kết luận ---
    slide = content_slide(prs, "Tổng kết", "Kết luận, giới hạn và hướng phát triển", 11)
    add_panel(slide, MARGIN, 1.55, 3.85, 4.85, "Đã đạt")
    add_bullets(
        slide,
        [
            "OOP kho có nhật ký",
            "NumPy · Pandas · Merge/Pivot",
            "10+ biểu đồ có nhận xét",
            "Baseline + 2 mô hình",
            "RFM + K-Means",
            "Pipeline chạy 1 lệnh",
        ],
        MARGIN + 0.28,
        2.2,
        3.3,
        13.5,
    )
    add_panel(slide, 4.78, 1.55, 3.85, 4.85, "Giới hạn")
    add_bullets(
        slide,
        [
            "Giao dịch synthetic",
            "FX cố định làm lệch doanh thu",
            "Chỉ 90 SP cho bài tồn kho",
            "Lớp low_stock mất cân bằng",
            "Chưa đánh giá theo thời gian",
        ],
        5.06,
        2.2,
        3.3,
        13.5,
    )
    add_panel(slide, 8.75, 1.55, 3.85, 4.85, "Hướng phát triển")
    add_bullets(
        slide,
        [
            "Dữ liệu bán hàng thật ẩn danh",
            "Dự báo nhu cầu theo sản phẩm",
            "Hiệu chỉnh ngưỡng reorder",
            "Theo dõi drift mô hình",
            "Kiểm thử và version dữ liệu",
        ],
        9.03,
        2.2,
        3.3,
        13.5,
    )

    # --- 12. Q&A ---
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, INK)
    add_rect(slide, 0, 0, SLIDE_W, 0.16, ACCENT)
    frame = add_text(slide, 1.5, 2.55, 10.3, 1.7)
    style_text(
        frame,
        [("CẢM ƠN THẦY/CÔ", 15, True, ACCENT), ("Q & A", 52, True, WHITE)],
        align=PP_ALIGN.CENTER,
        space_after=8,
    )
    add_rect(slide, 5.9, 4.45, 1.5, 0.045, ACCENT)
    frame = add_text(slide, 1.5, 4.75, 10.3, 1.1)
    style_text(
        frame,
        [
            ("Chạy lại pipeline:  python src/run_all.py", 14, False, PRGBColor(178, 195, 210)),
            ("Demo hệ thống:  python -m streamlit run src/app_dashboard.py", 14, False, PRGBColor(178, 195, 210)),
        ],
        align=PP_ALIGN.CENTER,
        space_after=5,
    )

    out = REPORTS / "SLIDE_CHUYEN_DE_6.pptx"
    prs.save(out)
    return out


def main() -> None:
    report = build_report()
    slides = build_slides()
    print(f"Created: {report}")
    print(f"Created: {slides}")


if __name__ == "__main__":
    main()
